#!/usr/bin/env python3
"""
ROS 2 Mini-AMR deck — v5, full design-system redesign.
Grid: 64px margins, 564/564 two-column split (24px gutter), shared content
top-line at y=180 across every slide. Type scale collapsed to 7 sizes.
Radius collapsed to 3 tiers. Icons standardized (dot=8, tick=20). B-roll
video embedded at its native 1002x652 aspect (zero letterbox/distortion).
"""
import os
from PIL import Image

from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement

D = "/home/nj/amr_ws/deliverables"
OUT = os.path.join(D, "ROS2_MiniAMR_Presentation.pptx")
VIDEO_MAIN = os.path.join(D, "mini_amr_final_demo.mp4")
CLIP_ROBOT, POSTER_ROBOT = os.path.join(D, "clip_robot.mp4"), os.path.join(D, "poster_robot.png")
CLIP_SAFETY, POSTER_SAFETY = os.path.join(D, "clip_safety.mp4"), os.path.join(D, "poster_safety.png")
CLIP_SLAM, POSTER_SLAM = os.path.join(D, "clip_slam.mp4"), os.path.join(D, "poster_slam.png")
CLIP_NAV, POSTER_NAV = os.path.join(D, "clip_nav.mp4"), os.path.join(D, "poster_nav.png")
HERO = os.path.join(D, "title_hero_new.png")
POSTER_MAIN = os.path.join(D, "nav_goal_new.png")

# ------------------------------------------------------------------ palette (unchanged)
DEEP, BLUE, CYAN = RGBColor(0x0F,0x4C,0x81), RGBColor(0x4A,0x90,0xE2), RGBColor(0x00,0xA6,0xD6)
INK, MUTED = RGBColor(0x22,0x22,0x22), RGBColor(0x5B,0x6B,0x7B)
PANEL, PANEL2, LINE = RGBColor(0xF3,0xF7,0xFB), RGBColor(0xE9,0xF1,0xF8), RGBColor(0xDB,0xE4,0xEC)
OK, WARN, WHITE = RGBColor(0x1E,0x9E,0x6A), RGBColor(0xC9,0x7A,0x16), RGBColor(0xFF,0xFF,0xFF)
STAGE, DANGER, SCANBG = RGBColor(0x0B,0x21,0x38), RGBColor(0xC0,0x39,0x2B), RGBColor(0x14,0x17,0x1B)
SANS, SEMI, MONO = "Segoe UI", "Segoe UI Semibold", "Consolas"

# ------------------------------------------------------------------ grid & type tokens
W, H = 1280, 720
MARGIN = 64
CONTENT_L, CONTENT_R = MARGIN, W - MARGIN            # 64, 1216
CONTENT_W = CONTENT_R - CONTENT_L                     # 1152
COL_GAP = 24
COL_W = (CONTENT_W - COL_GAP) / 2                     # 564
COL1_X, COL2_X = CONTENT_L, CONTENT_L + COL_W + COL_GAP  # 64, 652
CONTENT_TOP = 180
CONTENT_BOTTOM = 646
FOOTER_LINE_Y, FOOTER_TEXT_Y = 664, 674

T_EYEBROW, T_H1, T_H1_BIG, T_H2 = 13, 36, 40, 18
T_BODY, T_CAPTION, T_LABEL, T_MICRO, T_MONO, T_FOOTER = 14, 12, 11, 10, 12, 12

R_CARD, R_TAG, R_PILL = 0.06, 0.20, 0.5
ICON_DOT, ICON_TICK = 8, 20
BROLL_ASPECT = 1002 / 652     # 1.5368 — native clip aspect, zero letterbox
MASTER_ASPECT = 1920 / 1080   # 1.7778 — slide 8 only

EMU_PER_PX = 9525
def PX(v): return Emu(int(round(v * EMU_PER_PX)))
def PT(px): return Pt(px * 0.75)


def set_letter_spacing(run, hundredths_pt):
    run._r.get_or_add_rPr().set("spc", str(int(hundredths_pt)))


def no_autofit(tf):
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0


def add_text(slide, l, t, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             line_spacing=1.15, wrap=True):
    tb = slide.shapes.add_textbox(PX(l), PX(t), PX(w), PX(h))
    tf = tb.text_frame
    no_autofit(tf)
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    first = True
    for para in runs:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        p.line_spacing = line_spacing
        for item in para:
            text, size, color, bold, font = item[:5]
            spc = item[5] if len(item) > 5 else None
            r = p.add_run()
            r.text = text
            r.font.size = PT(size)
            r.font.color.rgb = color
            r.font.bold = bold
            r.font.name = font
            if spc:
                set_letter_spacing(r, spc)
    return tb


def simple_text(slide, l, t, w, h, text, size, color, bold=False, font=SANS,
                 align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.15):
    paras = [[(line, size, color, bold, font)] for line in text.split("\n")]
    return add_text(slide, l, t, w, h, paras, align=align, anchor=anchor, line_spacing=line_spacing)


def add_rect(slide, l, t, w, h, fill=None, line_color=None, line_w=1.0):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, PX(l), PX(t), PX(w), PX(h))
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line_color is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_color; shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    return shp


def add_card(slide, l, t, w, h, fill=PANEL, line_color=LINE, line_w=1.0, radius=R_CARD):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, PX(l), PX(t), PX(w), PX(h))
    shp.adjustments[0] = radius
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line_color is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_color; shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    return shp


def set_transparency(shape, alpha_pct):
    srgbClr = shape.fill.fore_color._xFill.find(qn("a:srgbClr"))
    a = OxmlElement("a:alpha"); a.set("val", str(int(alpha_pct * 1000)))
    srgbClr.append(a)


def add_line(slide, x1, y1, x2, y2, color=LINE, width=1.0, dash=None):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, PX(x1), PX(y1), PX(x2), PX(y2))
    conn.line.color.rgb = color; conn.line.width = Pt(width)
    if dash:
        d = OxmlElement("a:prstDash"); d.set("val", dash)
        conn.line._get_or_add_ln().append(d)
    return conn


def add_arrow(slide, x1, y1, x2, y2, color=BLUE, width=1.75, dash=None):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, PX(x1), PX(y1), PX(x2), PX(y2))
    conn.line.color.rgb = color; conn.line.width = Pt(width)
    ln = conn.line._get_or_add_ln()
    if dash:
        d = OxmlElement("a:prstDash"); d.set("val", dash); ln.append(d)
    te = OxmlElement("a:tailEnd"); te.set("type", "triangle"); te.set("w", "med"); te.set("len", "med")
    ln.append(te)
    return conn


def add_chip(slide, l, t, text, w=None, fill=PANEL2, tcolor=DEEP, size=T_LABEL,
             line_color=LINE, h=28, bold=True, font=SEMI):
    if w is None:
        w = 26 + len(text) * (size * 0.62)
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, PX(l), PX(t), PX(w), PX(h))
    shp.adjustments[0] = R_PILL
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = line_color; shp.line.width = Pt(1.0)
    shp.shadow.inherit = False
    tf = shp.text_frame
    no_autofit(tf); tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.size = PT(size); r.font.bold = bold; r.font.color.rgb = tcolor; r.font.name = font
    return shp, w


def add_dot(slide, l, t, color=CYAN, size=ICON_DOT):
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, PX(l), PX(t), PX(size), PX(size))
    dot.fill.solid(); dot.fill.fore_color.rgb = color
    dot.line.fill.background(); dot.shadow.inherit = False
    return dot


def add_tick(slide, l, t, ring_color=CYAN, mark_color=DEEP, size=ICON_TICK):
    circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, PX(l), PX(t), PX(size), PX(size))
    circ.fill.solid(); circ.fill.fore_color.rgb = PANEL
    circ.line.color.rgb = ring_color; circ.line.width = Pt(1.5)
    circ.shadow.inherit = False
    tf = circ.text_frame
    no_autofit(tf); tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = "✓"
    r.font.size = PT(size * 0.6); r.font.bold = True; r.font.color.rgb = mark_color; r.font.name = SANS
    return circ


def add_hdr(slide, num_label, title):
    add_text(slide, MARGIN, 52, 900, 22,
             [[("—  ", T_EYEBROW, CYAN, True, SEMI), (num_label, T_EYEBROW, DEEP, True, SEMI, 40)]],
             anchor=MSO_ANCHOR.MIDDLE)
    simple_text(slide, MARGIN, 82, CONTENT_W, 48, title, T_H1, DEEP, bold=True, font=SEMI)


def add_footer(slide, idx, total=10):
    add_line(slide, MARGIN, FOOTER_LINE_Y, W - MARGIN, FOOTER_LINE_Y, color=LINE, width=1.0)
    simple_text(slide, MARGIN, FOOTER_TEXT_Y, 400, 24, "ROS 2 Mini-AMR Platform", T_FOOTER, DEEP, bold=True, font=SEMI)
    simple_text(slide, W - MARGIN - 400, FOOTER_TEXT_Y, 400, 24, f"Jarupat Jaruvatee  ·  {idx:02d} / {total}",
                T_FOOTER, MUTED, font=SANS, align=PP_ALIGN.RIGHT)


def fit_contain(iw, ih, bw, bh):
    s = min(bw / iw, bh / ih)
    return iw * s, ih * s


def add_media_frame(slide, l, t, w, h, bg=SCANBG, border=LINE, radius_px=8):
    frame = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, PX(l), PX(t), PX(w), PX(h))
    frame.adjustments[0] = min(radius_px / min(w, h), 0.5)
    frame.fill.solid(); frame.fill.fore_color.rgb = bg
    frame.line.color.rgb = border; frame.line.width = Pt(1.0)
    frame.shadow.inherit = False
    return frame


def add_video_native(slide, video_path, poster_path, l, t, box_w, aspect, bg=SCANBG, border=LINE):
    """Place video at box_w wide, height derived from `aspect` (no distortion,
    no letterbox — the frame IS the video's own shape)."""
    box_h = box_w / aspect
    add_media_frame(slide, l, t, box_w, box_h, bg=bg, border=border)
    slide.shapes.add_movie(video_path, PX(l + 2), PX(t + 2), PX(box_w - 4), PX(box_h - 4),
                            poster_frame_image=poster_path, mime_type="video/mp4")
    return box_h


def add_image_contain(slide, path, l, t, w, h, bg=SCANBG, border=LINE, radius_px=8):
    add_media_frame(slide, l, t, w, h, bg=bg, border=border, radius_px=radius_px)
    iw, ih = Image.open(path).size
    fw, fh = fit_contain(iw, ih, w - 4, h - 4)
    slide.shapes.add_picture(path, PX(l + (w - fw) / 2), PX(t + (h - fh) / 2), PX(fw), PX(fh))


def crop_fracs_for_cover(iw, ih, bw, bh):
    img_ar, box_ar = iw / ih, bw / bh
    if img_ar > box_ar:
        new_w = ih * box_ar
        c = (iw - new_w) / 2 / iw
        return c, c, 0.0, 0.0
    new_h = iw / box_ar
    c = (ih - new_h) / 2 / ih
    return 0.0, 0.0, c, c


def add_image_cover(slide, path, l, t, w, h):
    iw, ih = Image.open(path).size
    cl, cr, ct, cb = crop_fracs_for_cover(iw, ih, w, h)
    pic = slide.shapes.add_picture(path, PX(l), PX(t), PX(w), PX(h))
    pic.crop_left, pic.crop_right, pic.crop_top, pic.crop_bottom = cl, cr, ct, cb
    return pic


def blank_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid(); s.background.fill.fore_color.rgb = WHITE
    return s


def wide_label(slide, x1, x2, y, text, color, min_w=170, h=16):
    mid = (x1 + x2) / 2
    w = max(min_w, x2 - x1)
    simple_text(slide, mid - w / 2, y, w, h, text, T_MICRO, color, bold=True, font=SEMI, align=PP_ALIGN.CENTER)


# ================================================================== deck
prs = Presentation()
prs.slide_width, prs.slide_height = PX(W), PX(H)

GOAL, REACHED = (1.00, 0.80), (1.004, 0.836)
DX, DY = REACHED[0] - GOAL[0], REACHED[1] - GOAL[1]
ERR_M = (DX ** 2 + DY ** 2) ** 0.5
ERR_CM = ERR_M * 100.0

NOTES = [
 "[~30 s] Good morning. I'm Jarupat Jaruvatee. In this internship I built a ROS 2-based Mini-AMR — an autonomous mobile robot developed entirely in simulation, from the robot model all the way to autonomous navigation.",
 "[~45 s] Autonomous mobile robots are becoming central to smart manufacturing and warehouse automation. Developing them directly on hardware is slow, costly, and risky — building the full stack in simulation first lets us design and tune everything safely. ROS 2 is the natural foundation.",
 "[~40 s] The objective was the complete AMR pipeline end to end: robot model, motion control and odometry, a LiDAR with perception and safety, SLAM mapping, and finally autonomous navigation with Nav2.",
 "[~60 s] Keyboard teleop publishes on /cmd_vel in MANUAL mode. The odometry node integrates that into the odom-to-base transform. During mapping, SLAM Toolbox produces the map; during AUTONOMOUS navigation, Nav2 localizes, plans, and drives. Manual and autonomous control never run at the same time.",
 "[~60 s] The robot is described in URDF/Xacro: a base_footprint and base_link chassis, four wheels, and a laser_link for the LiDAR. On the right is a live video, captured this session, of the robot driving with the LiDAR scan visible.",
 "[~50 s] A scan-analyzer node tracks the nearest obstacle; a safety-zone visualizer publishes RViz markers that turn red on a near obstacle. The video on the right is a fresh live capture showing the DANGER state.",
 "[~60 s] SLAM Toolbox performs online mapping — the video shows the occupancy grid being built live. Nav2 handles autonomous navigation: AMCL localizes, the planner produces a path, and the controller follows it.",
 f"[~60 s] The fully verified run: robot simulation, mapping, localization, and navigation all succeeded. The robot reached ({REACHED[0]:.3f}, {REACHED[1]:.3f}) against a goal of (1.00, 0.80) — error {ERR_CM:.1f} cm, BT: SUCCEEDED.",
 "[~50 s] The hardest problems were real engineering issues: a TF conflict, Nav2 costmap tuning, and activating the SLAM lifecycle node. Looking forward: multi-AMR fleets, a full warehouse simulation, and real-robot deployment.",
 f"[~30 s] A complete ROS 2 AMR simulation platform, reaching the goal within {ERR_CM:.1f} cm, fully verified with live ROS 2 runs. Thank you.",
]


def set_notes(slide, idx):
    slide.notes_slide.notes_text_frame.text = NOTES[idx]


# ------------------------------------------------------------------ SLIDE 1 — TITLE
s = blank_slide(prs)
RIGHT_W = 545
LEFT_W = W - RIGHT_W
add_image_cover(s, HERO, LEFT_W, 0, RIGHT_W, H)
tint = add_rect(s, LEFT_W, 0, RIGHT_W, H, fill=DEEP); set_transparency(tint, 10)
badge = add_card(s, LEFT_W + 22, H - 74, 470, 46, fill=STAGE, line_color=None, radius=0.3)
set_transparency(badge, 80)
add_text(s, LEFT_W + 38, H - 67, 440, 32,
         [[("Live capture · ", T_LABEL, WHITE, False, SANS), ("Nav2 goal reached", T_LABEL, CYAN, True, SEMI),
           (" in RViz2", T_LABEL, WHITE, False, SANS)]], anchor=MSO_ANCHOR.MIDDLE)

add_text(s, MARGIN, 54, 620, 22,
         [[("—  ", T_EYEBROW, CYAN, True, SEMI), ("INTERNSHIP FINAL PRESENTATION · ROBOTICS R&D", T_EYEBROW, DEEP, True, SEMI, 40)]],
         anchor=MSO_ANCHOR.MIDDLE)
simple_text(s, MARGIN, 86, 620, 100, "ROS 2-based Mini-AMR Simulation\n& Navigation Platform",
            T_H1_BIG, DEEP, bold=True, font=SEMI, line_spacing=1.1)
simple_text(s, MARGIN, 190, 620, 50,
            "A fully simulated autonomous mobile robot — modeling, motion control, "
            "sensing, SLAM mapping, and Nav2 autonomous navigation.",
            T_BODY + 1, MUTED, line_spacing=1.4)

meta = [("Presenter", "Jarupat Jaruvatee"), ("University", "King Mongkut's University of Technology Thonburi (KMUTT)"),
        ("Laboratory", "MEDAL Lab · National Taiwan University of Science and Technology (NTUST)"),
        ("Supervisor", "[Supervisor name]"), ("Date", "August 6, 2026")]
ROW_H = 36
y = 268
for label, val in meta:
    simple_text(s, MARGIN, y, 108, ROW_H, label, T_BODY, DEEP, bold=True, font=SEMI, anchor=MSO_ANCHOR.MIDDLE)
    tb = add_text(s, 182, y, 500, ROW_H, [[(val, T_BODY, INK, False, SANS)]], anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.2)
    if val.startswith("["):
        run = tb.text_frame.paragraphs[0].runs[0]
        run.font.italic = True; run.font.color.rgb = WARN
    y += ROW_H

chip_y = 470
cx = MARGIN
for label in ("ROS 2 Jazzy", "SLAM Toolbox", "Nav2", f"Goal error {ERR_CM:.1f} cm"):
    _, cw = add_chip(s, cx, chip_y, label, h=30)
    cx += cw + 10
set_notes(s, 0)

# ------------------------------------------------------------------ SLIDE 2 — BACKGROUND
s = blank_slide(prs)
add_hdr(s, "02 · BACKGROUND & MOTIVATION", "Why autonomous mobile robots — and why simulate first")

cards = [("Warehouse & smart manufacturing",
          "AMRs move material autonomously across dynamic factory and warehouse floors, "
          "replacing fixed conveyors and manual transport."),
         ("Simulation-first development",
          "Building and tuning the full stack in simulation de-risks the work — no hardware "
          "cost, damage, or safety risk while iterating."),
         ("ROS 2 as the foundation",
          "A modular, real-time-capable middleware (DDS) with a mature navigation ecosystem — "
          "the standard for research and industry robotics.")]
card_h, card_gap = 140, 27
cy = CONTENT_TOP
for head, body in cards:
    add_card(s, COL1_X, cy, COL_W, card_h)
    add_dot(s, COL1_X + 20, cy + 24)
    simple_text(s, COL1_X + 40, cy + 16, COL_W - 56, 26, head, T_H2, DEEP, bold=True, font=SEMI)
    simple_text(s, COL1_X + 40, cy + 50, COL_W - 56, 80, body, T_BODY, INK, line_spacing=1.45)
    cy += card_h + card_gap

diag_cx, diag_cy = COL2_X + COL_W / 2, CONTENT_TOP + (CONTENT_BOTTOM - CONTENT_TOP - 24) / 2
add_card(s, diag_cx - 78, diag_cy - 32, 156, 64, fill=DEEP, line_color=None, radius=R_PILL)
tf = s.shapes[-1].text_frame; no_autofit(tf); tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "ROS 2 (DDS)"; r.font.size = PT(T_H2 - 2); r.font.bold = True
r.font.color.rgb = WHITE; r.font.name = SEMI

spokes = [("Robot Model", -180, -130), ("Sensors", 180, -130), ("Simulation", -180, 130), ("Navigation", 180, 130)]
for label, dx, dy in spokes:
    add_line(s, diag_cx, diag_cy, diag_cx + dx, diag_cy + dy, color=BLUE, width=1.5)
for label, dx, dy in spokes:
    bw, bh = 152, 54
    sx, sy = diag_cx + dx, diag_cy + dy
    add_card(s, sx - bw / 2, sy - bh / 2, bw, bh, fill=PANEL2, line_color=CYAN, radius=R_TAG)
    tf = s.shapes[-1].text_frame; no_autofit(tf); tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = label; r.font.size = PT(T_BODY); r.font.bold = True
    r.font.color.rgb = DEEP; r.font.name = SEMI
simple_text(s, COL2_X, diag_cy + 170, COL_W, 24, "Nodes exchange data over topics; ROS 2 (DDS) is the shared middleware",
            T_CAPTION, MUTED, align=PP_ALIGN.CENTER)

add_footer(s, 2); set_notes(s, 1)

# ------------------------------------------------------------------ SLIDE 3 — OBJECTIVES
s = blank_slide(prs)
add_hdr(s, "03 · PROJECT OBJECTIVES", "A complete AMR development pipeline")

items = [[("Learn and apply ", False), ("ROS 2 architecture", True), (" — nodes, topics, TF, launch", False)],
         [("Build a ", False), ("URDF/Xacro robot model", True), (" with a full TF tree", False)],
         [("Implement ", False), ("motion control", True), (" and simulated odometry", False)],
         [("Simulate a ", False), ("LiDAR sensor", True), (" with perception & safety nodes", False)],
         [("Generate a map with ", False), ("SLAM Toolbox", True)],
         [("Achieve ", False), ("autonomous navigation", True), (" with Nav2", False)]]
item_h = ((CONTENT_BOTTOM - CONTENT_TOP) // len(items))
y = CONTENT_TOP + 6
for parts in items:
    add_tick(s, COL1_X, y)
    runs = [(t, T_BODY + 2, DEEP if b else INK, b, SEMI if b else SANS) for t, b in parts]
    add_text(s, COL1_X + 34, y - 1, COL_W - 34, ICON_TICK + 6, [runs], anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.2)
    y += item_h

steps = [("01", "Robot Model", "URDF/Xacro · TF"), ("02", "Control", "/cmd_vel · odometry"),
         ("03", "Sensors", "LiDAR /scan · safety"), ("04", "SLAM", "occupancy grid"),
         ("05", "Navigation", "Nav2 · goal reaching")]
fh, fgap = 64, 22
total_flow_h = len(steps) * fh + (len(steps) - 1) * fgap
fy = CONTENT_TOP + ((CONTENT_BOTTOM - CONTENT_TOP) - total_flow_h) / 2
for i, (n, head, sub) in enumerate(steps):
    yy = fy + i * (fh + fgap)
    add_card(s, COL2_X, yy, COL_W, fh, radius=R_TAG)
    add_rect(s, COL2_X, yy, 5, fh, fill=CYAN)
    simple_text(s, COL2_X + 22, yy + 9, 30, 22, n, T_CAPTION, BLUE, font=MONO)
    simple_text(s, COL2_X + 60, yy + 8, COL_W - 84, 24, head, T_H2 - 1, DEEP, bold=True, font=SEMI)
    simple_text(s, COL2_X + 60, yy + 34, COL_W - 84, 22, sub, T_CAPTION, MUTED)
    if i < len(steps) - 1:
        add_arrow(s, COL2_X + COL_W / 2, yy + fh, COL2_X + COL_W / 2, yy + fh + fgap, width=1.75)

add_footer(s, 3); set_notes(s, 2)

# ------------------------------------------------------------------ SLIDE 4 — ARCHITECTURE (full-width)
s = blank_slide(prs)
add_hdr(s, "04 · SYSTEM ARCHITECTURE", "ROS 2 node & data-flow pipeline")

lane1_y, lane2_y = CONTENT_TOP + 8, CONTENT_TOP + 8 + 182
simple_text(s, MARGIN, lane1_y - 24, 400, 20, "CONTROL & MODEL", T_CAPTION, MUTED, bold=True, font=SEMI)
simple_text(s, MARGIN, lane2_y - 24, 400, 20, "SENSING · MAPPING · AUTONOMY", T_CAPTION, MUTED, bold=True, font=SEMI)
add_line(s, MARGIN, lane1_y + 96, W - MARGIN, lane1_y + 96, color=LINE, dash="dash")

row1 = [("Keyboard", "teleop_twist", MARGIN, 168, PANEL, DEEP),
        ("Odometry / Mecanum", "fake_odom_publisher", MARGIN + 168 + 22, 198, PANEL, DEEP)]
boxh = 72
x = MARGIN
r1 = []
for head, sub, bw in [("Keyboard", "teleop_twist", 168), ("Odometry / Mecanum", "fake_odom_publisher", 198),
                       ("Robot Model", "URDF · robot_state", 168), ("TF + Odometry", "map→odom→base", 188)]:
    fill, border = (PANEL2, CYAN) if head == "TF + Odometry" else (PANEL, DEEP)
    add_card(s, x, lane1_y, bw, boxh, fill=fill, line_color=border, radius=R_TAG)
    simple_text(s, x + 10, lane1_y + 11, bw - 20, 22, head, T_BODY + 1, DEEP, bold=True, font=SEMI, align=PP_ALIGN.CENTER)
    simple_text(s, x + 10, lane1_y + 38, bw - 20, 20, sub, T_CAPTION - 1, MUTED, align=PP_ALIGN.CENTER)
    r1.append((x, x + bw)); x += bw + 22

pairs1 = [(r1[0][1], r1[1][0], "/cmd_vel", "MANUAL", BLUE), (r1[1][1], r1[2][0], None, None, BLUE),
          (r1[2][1], r1[3][0], "/odom", None, BLUE)]
for x1, x2, lbl, sub_lbl, col in pairs1:
    yc = lane1_y + boxh / 2
    add_arrow(s, x1, yc, x2, yc, color=col, width=1.75)
    if lbl:
        simple_text(s, x1, yc - 19, x2 - x1, 15, lbl, T_MONO - 1, CYAN, font=MONO, align=PP_ALIGN.CENTER)
    if sub_lbl:
        wide_label(s, x1, x2, lane1_y + boxh + 5, sub_lbl, BLUE)

x = MARGIN
r2 = []
row2 = [("LiDAR Sensor", "fake_scan_publisher", 168, PANEL, DEEP),
        ("SLAM Toolbox", "online mapping → /map", 198, PANEL, DEEP),
        ("Nav2 Navigation Stack", "AMCL · planner · controller · BT", 428, PANEL2, CYAN),
        ("Navigation", "Goal reached", 166, DEEP, None)]
for head, sub, bw, fill, border in row2:
    tcol = WHITE if fill == DEEP else DEEP
    subcol = RGBColor(0xbf, 0xe6, 0xf4) if fill == DEEP else MUTED
    add_card(s, x, lane2_y, bw, boxh, fill=fill, line_color=border, radius=R_TAG)
    simple_text(s, x + 10, lane2_y + 11, bw - 20, 22, head, T_BODY + 1, tcol, bold=True, font=SEMI, align=PP_ALIGN.CENTER)
    simple_text(s, x + 10, lane2_y + 38, bw - 20, 20, sub, T_CAPTION - 1, subcol, align=PP_ALIGN.CENTER)
    r2.append((x, x + bw)); x += bw + 22

pairs2 = [(r2[0][1], r2[1][0], "/scan", None), (r2[1][1], r2[2][0], "/map", None),
          (r2[2][1], r2[3][0], "/cmd_vel", "AUTONOMOUS")]
for x1, x2, lbl, sub_lbl in pairs2:
    yc = lane2_y + boxh / 2
    add_arrow(s, x1, yc, x2, yc, color=CYAN if sub_lbl else BLUE, width=2.0 if sub_lbl else 1.75)
    if lbl:
        simple_text(s, x1, yc - 19, x2 - x1, 15, lbl, T_MONO - 1, CYAN, font=MONO, align=PP_ALIGN.CENTER)
    if sub_lbl:
        wide_label(s, x1, x2, lane2_y + boxh + 5, sub_lbl, CYAN)

tf_cx = r1[3][0] + (r1[3][1] - r1[3][0]) / 2
slam_cx = r2[1][0] + (r2[1][1] - r2[1][0]) / 2
nav2_cx = r2[2][0] + (r2[2][1] - r2[2][0]) / 2
bus_y = lane1_y + boxh + 30
add_line(s, tf_cx, lane1_y + boxh, tf_cx, bus_y, color=BLUE, width=1.5, dash="dash")
add_line(s, slam_cx, bus_y, tf_cx, bus_y, color=BLUE, width=1.5, dash="dash")
add_arrow(s, slam_cx, bus_y, slam_cx, lane2_y, color=BLUE, width=1.5, dash="dash")
add_arrow(s, nav2_cx, bus_y, nav2_cx, lane2_y, color=BLUE, width=1.5, dash="dash")
simple_text(s, 610, bus_y - 24, 340, 16, "TF: map→odom→base_link→laser_link", T_MONO - 1, BLUE, font=MONO, align=PP_ALIGN.CENTER)

note_y = lane2_y + boxh + 40
add_card(s, MARGIN, note_y, CONTENT_W, 34, fill=PANEL2, radius=0.25)
add_text(s, MARGIN + 20, note_y + 6, CONTENT_W - 40, 22,
         [[("Note — ", T_CAPTION + 1, DEEP, True, SEMI),
           ("Manual teleoperation and autonomous navigation run separately, so /cmd_vel is never published by both sources at once.",
            T_CAPTION + 1, DEEP, False, SANS)]], anchor=MSO_ANCHOR.MIDDLE)

add_footer(s, 4); set_notes(s, 3)

# ------------------------------------------------------------------ SLIDE 5 — ROBOT MODELING (video)
s = blank_slide(prs)
add_hdr(s, "05 · ROBOT MODELING & CONTROL", "URDF/Xacro model and the TF tree")

card_h1, card_h2, gap = 96, 108, 16
add_card(s, COL1_X, CONTENT_TOP, COL_W, card_h1)
add_dot(s, COL1_X + 20, CONTENT_TOP + 20)
simple_text(s, COL1_X + 40, CONTENT_TOP + 12, COL_W - 56, 24, "Robot description", T_H2, DEEP, bold=True, font=SEMI)
add_text(s, COL1_X + 40, CONTENT_TOP + 42, COL_W - 56, 48,
         [[("base_footprint → base_link", T_BODY - 1, DEEP, True, MONO), (" (0.8 × 0.6 × 0.2 m body), four ", T_BODY - 1, INK, False, SANS),
           ("wheels", T_BODY - 1, DEEP, True, SEMI), (" (r = 0.1 m), and a ", T_BODY - 1, INK, False, SANS),
           ("laser_link", T_BODY - 1, DEEP, True, MONO), (" LiDAR mount.", T_BODY - 1, INK, False, SANS)]], line_spacing=1.35)

y2 = CONTENT_TOP + card_h1 + gap
add_card(s, COL1_X, y2, COL_W, card_h2)
add_dot(s, COL1_X + 20, y2 + 20)
simple_text(s, COL1_X + 40, y2 + 12, COL_W - 56, 24, "Motion & kinematics", T_H2, DEEP, bold=True, font=SEMI)
add_text(s, COL1_X + 40, y2 + 42, COL_W - 56, 62,
         [[("fake_odom_publisher", T_BODY - 1, DEEP, True, MONO), (" integrates ", T_BODY - 1, INK, False, SANS),
           ("/cmd_vel", T_BODY - 1, DEEP, True, MONO), (" into odometry + the moving ", T_BODY - 1, INK, False, SANS),
           ("odom→base_footprint", T_BODY - 1, DEEP, True, MONO),
           (" TF. A mecanum inverse-kinematics node maps body velocity to four wheel speeds.", T_BODY - 1, INK, False, SANS)]],
         line_spacing=1.35)

tfy = y2 + card_h2 + 42
tf_defs = [("map", 60, DEEP), ("odom", 60, RGBColor(0x12, 0x3f, 0x6b)), ("base_footprint", 140, CYAN), ("base_link", 110, BLUE)]
xs, x = [], COL1_X
for label, bw, col in tf_defs:
    add_card(s, x, tfy, bw, 38, fill=col, line_color=None, radius=R_TAG)
    tf_ = s.shapes[-1].text_frame; no_autofit(tf_); tf_.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf_.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = label; r.font.size = PT(T_CAPTION); r.font.bold = True
    r.font.color.rgb = WHITE; r.font.name = MONO
    xs.append((x, x + bw)); x += bw + 18
for i in range(len(xs) - 1):
    yc = tfy + 19
    add_arrow(s, xs[i][1], yc, xs[i + 1][0], yc, color=MUTED, width=1.4)
child_x = xs[-1][1] + 30
for i, lbl in enumerate(["4× wheel", "laser_link"]):
    cy_ = tfy - 19 + i * 30
    add_card(s, child_x, cy_, COL1_X + COL_W - child_x, 22, fill=PANEL2, line_color=BLUE, radius=R_TAG)
    tf_ = s.shapes[-1].text_frame; no_autofit(tf_); tf_.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf_.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = lbl; r.font.size = PT(T_CAPTION - 1); r.font.bold = True
    r.font.color.rgb = DEEP; r.font.name = MONO
    add_line(s, xs[-1][1], tfy + 19, child_x, cy_ + 11, color=BLUE, width=1.2)

vid_h = add_video_native(s, CLIP_ROBOT, POSTER_ROBOT, COL2_X, CONTENT_TOP, COL_W, BROLL_ASPECT)
add_chip(s, COL2_X, CONTENT_TOP + vid_h + 12, "▶ LIVE VIDEO", fill=CYAN, tcolor=WHITE, line_color=CYAN, h=22)
simple_text(s, COL2_X, CONTENT_TOP + vid_h + 44, COL_W, 36,
            "RViz2 recording — URDF model driving with the live LaserScan", T_CAPTION, MUTED,
            align=PP_ALIGN.CENTER, line_spacing=1.3)
legend_y = CONTENT_TOP + vid_h + 82
lx = COL2_X + (COL_W - 260) / 2
add_rect(s, lx, legend_y, 12, 12, fill=RGBColor(0x1f, 0x5f, 0xd0))
simple_text(s, lx + 18, legend_y - 4, 100, 20, "base_link", T_CAPTION, INK, font=MONO)
lx += 130
dot = s.shapes.add_shape(MSO_SHAPE.OVAL, PX(lx), PX(legend_y), PX(12), PX(12))
dot.fill.solid(); dot.fill.fore_color.rgb = DANGER; dot.line.fill.background(); dot.shadow.inherit = False
simple_text(s, lx + 18, legend_y - 4, 100, 20, "laser_link", T_CAPTION, INK, font=MONO)

add_footer(s, 5); set_notes(s, 4)

# ------------------------------------------------------------------ SLIDE 6 — SENSOR SIMULATION (video)
s = blank_slide(prs)
add_hdr(s, "06 · SENSOR SIMULATION", "Simulated LiDAR, perception & safety zones")

cards6 = [("Fake LaserScan publisher",
           [("Publishes ", False), ("/scan", True), (" in ", False), ("laser_link", True),
            (" — 360 rays at 5 Hz, 0.1–3.0 m range. Supports a clear (open-space) mode for stable localization.", False)]),
          ("Scan analyzer",
           [("Splits the scan into front / left / right sectors, tracks the nearest return, and raises ", False),
            ("/obstacle_alert", True), (" when an obstacle enters the safety threshold.", False)]),
          ("Safety-zone markers",
           [("Publishes RViz ", False), ("MarkerArray", True),
            (" zones that turn red on a near obstacle and stay blue when clear.", False)])]
card_h, card_gap = 122, 20
total_h = 3 * card_h + 2 * card_gap
cy = CONTENT_TOP + ((CONTENT_BOTTOM - CONTENT_TOP) - total_h) / 2
for head, parts in cards6:
    add_card(s, COL1_X, cy, COL_W, card_h)
    add_dot(s, COL1_X + 20, cy + 24)
    simple_text(s, COL1_X + 40, cy + 16, COL_W - 56, 26, head, T_H2, DEEP, bold=True, font=SEMI)
    runs = [(t, T_BODY, DEEP if mono else INK, mono, MONO if mono else SANS) for t, mono in parts]
    add_text(s, COL1_X + 40, cy + 48, COL_W - 56, 66, [runs], line_spacing=1.4)
    cy += card_h + card_gap

vid_h = add_video_native(s, CLIP_SAFETY, POSTER_SAFETY, COL2_X, CONTENT_TOP, COL_W, BROLL_ASPECT)
badge, bw = add_chip(s, COL2_X, CONTENT_TOP + vid_h + 12, "▶ LIVE VIDEO", fill=CYAN, tcolor=WHITE, line_color=CYAN, h=24)
add_text(s, COL2_X + bw + 10, CONTENT_TOP + vid_h + 12, COL_W - bw - 10, 40,
         [[("Safety zones from LiDAR ", T_CAPTION, INK, False, SANS), ("/scan", T_CAPTION, DEEP, True, MONO),
           (", drawn as a ", T_CAPTION, INK, False, SANS), ("MarkerArray", T_CAPTION, DEEP, True, MONO),
           (" — recorded in the DANGER state", T_CAPTION, INK, False, SANS)]], line_spacing=1.35)

legend_y = CONTENT_TOP + vid_h + 12 + 40 + 14
legend_h = CONTENT_BOTTOM - legend_y
add_card(s, COL2_X, legend_y, COL_W, legend_h, fill=PANEL2, radius=0.3)
zones = [("FRONT", DANGER), ("LEFT", BLUE), ("RIGHT", BLUE)]
zw = COL_W / len(zones)
for i, (label, col) in enumerate(zones):
    zx = COL2_X + i * zw
    sw_ = s.shapes.add_shape(MSO_SHAPE.OVAL, PX(zx + zw / 2 - 34), PX(legend_y + legend_h / 2 - 5), PX(10), PX(10))
    sw_.fill.solid(); sw_.fill.fore_color.rgb = col; sw_.line.fill.background(); sw_.shadow.inherit = False
    simple_text(s, zx + zw / 2 - 18, legend_y + legend_h / 2 - 9, 60, 18, label, T_MICRO, DEEP, bold=True, font=SEMI, anchor=MSO_ANCHOR.MIDDLE)

add_footer(s, 6); set_notes(s, 5)

# ------------------------------------------------------------------ SLIDE 7 — SLAM + NAV2 (both video)
s = blank_slide(prs)
add_hdr(s, "07 · SLAM & AUTONOMOUS NAVIGATION", "Mapping the space, then navigating it")

py, ph = CONTENT_TOP, CONTENT_BOTTOM - CONTENT_TOP
add_card(s, COL1_X, py, COL_W, ph)
add_card(s, COL2_X, py, COL_W, ph)

pad = 20
add_dot(s, COL1_X + pad, py + 26, color=BLUE)
simple_text(s, COL1_X + pad + 20, py + 18, 240, 26, "SLAM Toolbox", T_H2, DEEP, bold=True, font=SEMI)
add_chip(s, COL1_X + COL_W - pad - 100, py + 16, "MAPPING", w=100, fill=BLUE, tcolor=WHITE, line_color=BLUE, h=24)
sy = py + 58
for txt in ["Online async mapping from LiDAR", "Occupancy grid built while driving", "Lifecycle-managed · publishes map→odom"]:
    add_tick(s, COL1_X + pad, sy, ring_color=BLUE, size=18)
    simple_text(s, COL1_X + pad + 28, sy - 1, COL_W - 2 * pad - 28, 24, txt, T_BODY - 1, INK, anchor=MSO_ANCHOR.MIDDLE)
    sy += 28
vy = sy + 10
vid_w = COL_W - 2 * pad
vid_h1 = min(vid_w / BROLL_ASPECT, py + ph - 40 - vy)
add_video_native(s, CLIP_SLAM, POSTER_SLAM, COL1_X + pad, vy, vid_w, BROLL_ASPECT)
cap_y = vy + vid_h1 + 8
badge, bws = add_chip(s, COL1_X + pad, cap_y, "▶ LIVE", fill=BLUE, tcolor=WHITE, line_color=BLUE, size=T_MICRO, h=18)
simple_text(s, COL1_X + pad + bws + 8, cap_y, COL_W - 2 * pad - bws - 8, 20, "Occupancy grid recorded live this session", T_MICRO + 1, MUTED, anchor=MSO_ANCHOR.MIDDLE)

add_dot(s, COL2_X + pad, py + 26, color=CYAN)
simple_text(s, COL2_X + pad + 20, py + 18, 240, 26, "Nav2", T_H2, DEEP, bold=True, font=SEMI)
add_chip(s, COL2_X + COL_W - pad - 118, py + 16, "NAVIGATION", w=118, fill=CYAN, tcolor=WHITE, line_color=CYAN, h=24)
sy = py + 58
for txt in ["AMCL localization on the saved map", "NavFn global planner + MPPI controller", "Static + inflation costmaps, BT navigator"]:
    add_tick(s, COL2_X + pad, sy, ring_color=CYAN, size=18)
    simple_text(s, COL2_X + pad + 28, sy - 1, COL_W - 2 * pad - 28, 24, txt, T_BODY - 1, INK, anchor=MSO_ANCHOR.MIDDLE)
    sy += 28
vy = sy + 10
vid_h2 = min(vid_w / BROLL_ASPECT, py + ph - 46 - vy)
add_video_native(s, CLIP_NAV, POSTER_NAV, COL2_X + pad, vy, vid_w, BROLL_ASPECT)
cap_y = vy + vid_h2 + 8
badge, bw3 = add_chip(s, COL2_X + pad, cap_y, "▶ LIVE VIDEO", fill=CYAN, tcolor=WHITE, line_color=CYAN, size=T_MICRO, h=20)
simple_text(s, COL2_X + pad + bw3 + 8, cap_y, COL_W - 2 * pad - bw3 - 8, 20, "Global path → robot reaches the goal", T_MICRO + 1, MUTED, anchor=MSO_ANCHOR.MIDDLE)

add_footer(s, 7); set_notes(s, 6)

# ------------------------------------------------------------------ SLIDE 8 — RESULTS
s = blank_slide(prs)
add_hdr(s, "08 · EXPERIMENTAL RESULTS", "Verified end-to-end navigation")

metrics = [("Robot simulation", "Successful"), ("SLAM mapping", "Completed"),
           ("Localization", "Stable"), ("Navigation", "Successful")]
mgap = 16
mw = (COL_W - mgap) / 2
mh = 54
for i, (k, v) in enumerate(metrics):
    x = COL1_X + (i % 2) * (mw + mgap)
    y = CONTENT_TOP + (i // 2) * (mh + 14)
    add_card(s, x, y, mw, mh)
    add_text(s, x + 16, y, mw - 110, mh, [[(k, T_BODY, INK, False, SANS)]], anchor=MSO_ANCHOR.MIDDLE)
    add_tick(s, x + mw - 92, y + mh / 2 - 9, ring_color=OK, mark_color=OK, size=18)
    add_text(s, x + mw - 68, y, 60, mh, [[(v, T_CAPTION, OK, True, SEMI)]], anchor=MSO_ANCHOR.MIDDLE)

hy = CONTENT_TOP + 2 * (mh + 14) + 12
hero_h = 172
add_card(s, COL1_X, hy, COL_W, hero_h, fill=DEEP, line_color=None, radius=R_CARD)
add_text(s, COL1_X + 24, hy + 16, 260, 60, [[(f"{ERR_CM:.1f}", 52, WHITE, True, SEMI), (" cm", 20, WHITE, False, SEMI)]])
simple_text(s, COL1_X + 24, hy + 78, COL_W - 48, 20, "Final goal-position error (Euclidean)", T_CAPTION + 1, WHITE)
add_text(s, COL1_X + 24, hy + 102, COL_W - 48, 52,
         [[(f"goal ({GOAL[0]:.2f}, {GOAL[1]:.2f}) → reached ({REACHED[0]:.3f}, {REACHED[1]:.3f}) m",
            T_CAPTION - 1, RGBColor(0xbf, 0xe6, 0xf4), False, MONO)],
          [(f"Δx={DX:+.3f} m · Δy={DY:+.3f} m · err≈{ERR_M:.3f} m · BT: SUCCEEDED",
            T_CAPTION - 1, RGBColor(0xbf, 0xe6, 0xf4), False, MONO)]], line_spacing=1.5)

foot_y = hy + hero_h + 14
simple_text(s, COL1_X, foot_y, COL_W, 40,
            "Orientation error was monitored throughout the run. See the embedded "
            "recording (right) for the complete navigation task.", T_CAPTION, MUTED, line_spacing=1.4)

vid_w = COL_W
vid_h = vid_w / MASTER_ASPECT
vid_y = CONTENT_TOP + ((CONTENT_BOTTOM - CONTENT_TOP) - (vid_h + 34)) / 2
add_video_native(s, VIDEO_MAIN, POSTER_MAIN, COL2_X, vid_y, vid_w, MASTER_ASPECT)
simple_text(s, COL2_X, vid_y + vid_h + 10, vid_w, 22, "▶ Embedded demo video (60 s) — click to play", T_CAPTION, MUTED, align=PP_ALIGN.CENTER)

add_footer(s, 8); set_notes(s, 7)

# ------------------------------------------------------------------ SLIDE 9 — CHALLENGES / FUTURE
s = blank_slide(prs)
add_hdr(s, "09 · CHALLENGES & FUTURE WORK", "What was hard — and what comes next")

py, ph = CONTENT_TOP, CONTENT_BOTTOM - CONTENT_TOP
add_card(s, COL1_X, py, COL_W, ph)
add_card(s, COL2_X, py, COL_W, ph)

pad = 20
add_dot(s, COL1_X + pad, py + 24, color=WARN)
simple_text(s, COL1_X + pad + 20, py + 16, 300, 26, "Challenges solved", T_H2, WARN, bold=True, font=SEMI)
chal = [[("TF configuration", True), (" — resolved a fixed vs. moving odom→base_link conflict that pinned the robot", False)],
        [("Navigation tuning", True), (" — costmaps, inflation & goal tolerance for reliable planning", False)],
        [("Sensor integration", True), (" — activating the SLAM lifecycle node so mapping starts", False)],
        [("ROS 2 debugging", True), (" — lifecycle, QoS & TF timing across nodes", False)]]
row_h = (ph - 62) / len(chal)
sy = py + 62
for parts in chal:
    add_tick(s, COL1_X + pad, sy, ring_color=WARN, mark_color=WARN)
    runs = [(t, T_BODY, DEEP if b else INK, b, SEMI if b else SANS) for t, b in parts]
    add_text(s, COL1_X + pad + 32, sy - 1, COL_W - 2 * pad - 32, row_h - 10, [runs], line_spacing=1.3)
    sy += row_h

add_dot(s, COL2_X + pad, py + 24, color=CYAN)
simple_text(s, COL2_X + pad + 20, py + 16, 300, 26, "Future work", T_H2, DEEP, bold=True, font=SEMI)
fut = [[("Multi-AMR fleet management", True), (" & coordination", False)],
       [("Full warehouse simulation", True), (" environment", False)],
       [("Real-robot", True), (" deployment of the same stack", False)],
       [("Task scheduling", True), (" & higher-level mission control", False)]]
sy = py + 62
for parts in fut:
    add_tick(s, COL2_X + pad, sy)
    runs = [(t, T_BODY, DEEP if b else INK, b, SEMI if b else SANS) for t, b in parts]
    add_text(s, COL2_X + pad + 32, sy - 1, COL_W - 2 * pad - 32, row_h - 10, [runs], line_spacing=1.3)
    sy += row_h

add_footer(s, 9); set_notes(s, 8)

# ------------------------------------------------------------------ SLIDE 10 — CONCLUSION
s = blank_slide(prs)
add_hdr(s, "10 · CONCLUSION", "A complete, verified AMR simulation platform")

chip_labels = ["Robot modeling", "LiDAR integration", "SLAM mapping", "Nav2 navigation"]
widths = [26 + len(lbl) * (T_LABEL * 0.62) for lbl in chip_labels]
total_w = sum(widths) + 12 * (len(widths) - 1)
tx = W / 2 - total_w / 2
chip_y = CONTENT_TOP + 30
for label, cw in zip(chip_labels, widths):
    add_chip(s, tx, chip_y, label, w=cw, h=30)
    tx += cw + 12

simple_text(s, W / 2 - 420, chip_y + 54, 840, 60,
            f"Built on ROS 2 Jazzy — robot modeling, motion control, simulated sensing, "
            f"SLAM mapping, and autonomous navigation reaching the goal within {ERR_CM:.1f} cm.",
            T_BODY + 2, MUTED, align=PP_ALIGN.CENTER, line_spacing=1.4)

mid_y = (CONTENT_TOP + CONTENT_BOTTOM) / 2 + 30
simple_text(s, 0, mid_y, W, 76, "Thank You", 54, DEEP, bold=True, font=SEMI, align=PP_ALIGN.CENTER)
simple_text(s, 0, mid_y + 78, W, 40, "Questions?", 24, CYAN, bold=True, font=SEMI, align=PP_ALIGN.CENTER)
simple_text(s, 0, mid_y + 130, W, 26, "Jarupat Jaruvatee  ·  ROS 2 Mini-AMR Simulation & Navigation Platform",
            T_BODY, MUTED, align=PP_ALIGN.CENTER)

add_footer(s, 10); set_notes(s, 9)

# ------------------------------------------------------------------ save
prs.save(OUT)
print("saved", OUT, os.path.getsize(OUT) // 1024, "KB")
