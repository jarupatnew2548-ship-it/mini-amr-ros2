from pptx import Presentation
from pptx.util import Emu

PATH = "/home/nj/amr_ws/deliverables/ROS2_MiniAMR_Presentation.pptx"
EMU_PER_PX = 9525
W, H = 1280, 720

prs = Presentation(PATH)
issues = []

for si, slide in enumerate(prs.slides, start=1):
    boxes = []
    for shp in slide.shapes:
        try:
            l = shp.left / EMU_PER_PX
            t = shp.top / EMU_PER_PX
            w = shp.width / EMU_PER_PX
            h = shp.height / EMU_PER_PX
        except TypeError:
            continue
        r, b = l + w, t + h
        if l < -1 or t < -1 or r > W + 1 or b > H + 1:
            issues.append(f"S{si}: OUT OF BOUNDS {shp.shape_type} '{getattr(shp,'name','')}' "
                           f"l={l:.0f} t={t:.0f} r={r:.0f} b={b:.0f}")
        if shp.has_text_frame and shp.text_frame.text.strip():
            boxes.append((l, t, r, b, shp.text_frame.text.strip()[:30]))

    for i in range(len(boxes)):
        for j in range(i + 1, len(boxes)):
            l1, t1, r1, b1, txt1 = boxes[i]
            l2, t2, r2, b2, txt2 = boxes[j]
            ix = max(0, min(r1, r2) - max(l1, l2))
            iy = max(0, min(b1, b2) - max(t1, t2))
            area = ix * iy
            a1 = (r1 - l1) * (b1 - t1)
            a2 = (r2 - l2) * (b2 - t2)
            if area > 0.35 * min(a1, a2) and area > 200:
                issues.append(f"S{si}: OVERLAP '{txt1}' <-> '{txt2}' area={area:.0f}")

media_parts = [p.partname for p in prs.part.package.iter_parts() if "media" in str(p.partname)]
print(f"Slides: {len(prs.slides.__iter__.__self__._sldIdLst)}")
print(f"Media parts embedded: {len(media_parts)}")
for m in media_parts:
    print(" ", m)

print()
if issues:
    print(f"{len(issues)} ISSUES:")
    for i in issues:
        print(" -", i)
else:
    print("0 issues: no out-of-bounds shapes, no significant text-box overlaps.")
